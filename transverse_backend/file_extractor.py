#!/usr/bin/env python3
"""
File extraction service for various document formats.
Supports: PDF, XPS, EPUB, MOBI, FB2, CBZ, SVG, TXT, Images, DOCX, XLSX, PPTX, HWP
"""

import os
import zipfile
import tempfile
from pathlib import Path
import mimetypes

class FileExtractor:
    def __init__(self):
        self.supported_extensions = {
            '.pdf', '.xps', '.epub', '.mobi', '.fb2', '.cbz', 
            '.svg', '.txt', '.docx', '.xlsx', '.pptx', '.hwp',
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'
        }
        
    def is_supported(self, filename):
        """Check if file extension is supported."""
        ext = Path(filename).suffix.lower()
        return ext in self.supported_extensions
    
    def extract_text(self, file_path):
        """
        Extract text from uploaded file based on its type.
        
        Args:
            file_path (str): Path to the uploaded file
            
        Returns:
            dict: Contains extracted text, file info, and any errors
        """
        try:
            file_path = Path(file_path)
            extension = file_path.suffix.lower()
            
            result = {
                'filename': file_path.name,
                'extension': extension,
                'size': file_path.stat().st_size,
                'text': '',
                'error': None,
                'pages': 1,
                'metadata': {}
            }
            
            # Route to appropriate extraction method
            if extension == '.pdf':
                result.update(self._extract_pdf(file_path))
            elif extension == '.txt':
                result.update(self._extract_txt(file_path))
            elif extension == '.docx':
                result.update(self._extract_docx(file_path))
            elif extension == '.xlsx':
                result.update(self._extract_xlsx(file_path))
            elif extension == '.pptx':
                result.update(self._extract_pptx(file_path))
            elif extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
                result.update(self._extract_image(file_path))
            elif extension == '.epub':
                result.update(self._extract_epub(file_path))
            elif extension == '.svg':
                result.update(self._extract_svg(file_path))
            elif extension == '.cbz':
                result.update(self._extract_cbz(file_path))
            elif extension in ['.xps', '.mobi', '.fb2', '.hwp']:
                result.update(self._extract_unsupported(file_path, extension))
            else:
                result['error'] = f"Unsupported file type: {extension}"
                
            return result
            
        except Exception as e:
            return {
                'filename': str(file_path),
                'extension': extension if 'extension' in locals() else 'unknown',
                'size': 0,
                'text': '',
                'error': f"Error processing file: {str(e)}",
                'pages': 0,
                'metadata': {}
            }
    
    def _extract_pdf(self, file_path):
        """Extract text from PDF files with enhanced structure preservation."""
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            
            text_pages = []
            page_info = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # Extract text with structure using dict format for better formatting
                page_dict = page.get_text("dict")
                page_text = self._extract_structured_text(page_dict)
                text_pages.append(page_text)
                
                # Store page info for later use
                page_info.append({
                    'page_number': page_num + 1,
                    'width': page.rect.width,
                    'height': page.rect.height,
                    'rotation': page.rotation
                })
            
            doc.close()
            
            return {
                'text': '\n\n--- Page Break ---\n\n'.join(text_pages),
                'pages': len(text_pages),
                'metadata': {
                    'total_pages': len(text_pages),
                    'page_info': page_info,
                    'document_path': str(file_path)
                }
            }
        except ImportError:
            return {'error': 'PyMuPDF not installed. Run: pip install PyMuPDF'}
        except Exception as e:
            return {'error': f'PDF extraction error: {str(e)}'}
    
    def _extract_structured_text(self, page_dict):
        """Extract text while preserving structure from PyMuPDF dict format."""
        text_blocks = []
        
        for block in page_dict.get("blocks", []):
            if "lines" in block:  # Text block
                block_text = []
                for line in block["lines"]:
                    line_text = []
                    for span in line["spans"]:
                        if span["text"].strip():
                            line_text.append(span["text"])
                    if line_text:
                        block_text.append("".join(line_text))
                
                if block_text:
                    text_blocks.append("\n".join(block_text))
        
        return "\n\n".join(text_blocks)
    
    def extract_pdf_pages(self, file_path, page_numbers=None):
        """Extract text from specific PDF pages."""
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            
            if page_numbers is None:
                page_numbers = list(range(len(doc)))
            else:
                # Convert to 0-based indexing and validate
                page_numbers = [p-1 for p in page_numbers if 0 < p <= len(doc)]
            
            extracted_pages = []
            for page_num in page_numbers:
                page = doc[page_num]
                page_dict = page.get_text("dict")
                page_text = self._extract_structured_text(page_dict)
                
                extracted_pages.append({
                    'page_number': page_num + 1,
                    'text': page_text,
                    'width': page.rect.width,
                    'height': page.rect.height,
                    'rotation': page.rotation
                })
            
            doc.close()
            
            return {
                'success': True,
                'pages': extracted_pages,
                'total_pages': len(doc),
                'selected_pages': len(extracted_pages)
            }
            
        except ImportError:
            return {'success': False, 'error': 'PyMuPDF not installed'}
        except Exception as e:
            return {'success': False, 'error': f'PDF page extraction error: {str(e)}'}
    
    def translate_pdf_pages(self, file_path, page_numbers, target_language, translation_service):
        """Translate specific PDF pages and return a new PDF with translated text."""
        try:
            # Use the advanced PDF translator for better results
            from .pdf_translator_advanced import advanced_pdf_translator
            
            result = advanced_pdf_translator.translate_pdf_with_redaction(
                file_path, page_numbers, target_language, translation_service
            )
            
            return result
            
        except ImportError:
            # Fallback to basic method if advanced translator is not available
            return self._translate_pdf_basic(file_path, page_numbers, target_language, translation_service)
        except Exception as e:
            return {'success': False, 'error': f'PDF translation error: {str(e)}'}
    
    def _translate_pdf_basic(self, file_path, page_numbers, target_language, translation_service):
        """Basic PDF translation method as fallback."""
        try:
            import fitz  # PyMuPDF
            
            # Open original document
            original_doc = fitz.open(file_path)
            
            if page_numbers is None:
                page_numbers = list(range(1, len(original_doc) + 1))
            
            # Create new document for translated content
            translated_doc = fitz.open()
            
            for page_idx in page_numbers:
                if page_idx < 1 or page_idx > len(original_doc):
                    continue
                    
                # Get original page (convert to 0-based)
                original_page = original_doc[page_idx - 1]
                
                # Create new page with same dimensions
                new_page = translated_doc.new_page(width=original_page.rect.width, 
                                                 height=original_page.rect.height)
                
                # First, copy all images from original page
                self._copy_page_images(original_page, new_page)
                
                # Extract text blocks with positioning using redaction approach
                page_dict = original_page.get_text("dict")
                
                # Process each text block for translation
                for block in page_dict.get("blocks", []):
                    if "lines" in block:  # Text block
                        # Collect all text and formatting info from the block
                        block_text_parts = []
                        block_bbox = None
                        font_info = None
                        
                        # Get text and bounding box for the entire block
                        for line in block["lines"]:
                            for span in line["spans"]:
                                if span["text"].strip():
                                    block_text_parts.append(span["text"])
                                    
                                    # Get font information from first meaningful span
                                    if font_info is None:
                                        font_info = {
                                            'size': span.get('size', 12),
                                            'flags': span.get('flags', 0),
                                            'color': span.get('color', 0),
                                            'font': span.get('font', 'Times-Roman')
                                        }
                                    
                                    # Calculate bounding box for entire block
                                    bbox = fitz.Rect(span['bbox'])
                                    if block_bbox is None:
                                        block_bbox = bbox
                                    else:
                                        block_bbox = block_bbox | bbox
                        
                        if block_text_parts and block_bbox:
                            # Combine text and translate
                            original_text = " ".join(block_text_parts).strip()
                            
                            try:
                                translated_text = translation_service.translate(original_text, target_language)
                                if not translated_text or translated_text.strip() == "":
                                    translated_text = original_text  # Fallback
                            except Exception as e:
                                print(f"Translation error: {e}")
                                translated_text = original_text  # Fallback to original
                            
                            # Choose appropriate font based on target language and content
                            fontname = self._get_appropriate_font(target_language, font_info)
                            
                            # Determine text styling
                            is_bold = bool(font_info['flags'] & 2**4)
                            is_italic = bool(font_info['flags'] & 2**6)
                            
                            # Convert color from int to RGB tuple
                            color = self._convert_color(font_info['color'])
                            
                            # Insert translated text with better error handling
                            try:
                                # Try insert_htmlbox first for better Unicode support
                                font_size = max(8, min(font_info['size'], 24))  # Reasonable font size range
                                
                                # Create HTML with proper encoding
                                html_content = self._create_html_content(translated_text, font_size, color, is_bold, is_italic)
                                
                                # Try HTML insertion first (better for Unicode)
                                try:
                                    insertion_result = new_page.insert_htmlbox(block_bbox, html_content)
                                    if insertion_result[0] < 0:  # HTML insertion failed
                                        raise Exception("HTML insertion failed")
                                except:
                                    # Fallback to textbox with Unicode font
                                    insertion_result = new_page.insert_textbox(
                                        block_bbox,
                                        translated_text,
                                        fontname=fontname,
                                        fontsize=font_size,
                                        color=color,
                                        align=0,
                                        encoding=0  # Use default encoding for better compatibility
                                    )
                                    
                                    # If still failing, try with smaller font
                                    if insertion_result < 0:
                                        for smaller_size in [font_size * 0.8, font_size * 0.6, font_size * 0.4]:
                                            if smaller_size < 6:
                                                break
                                            insertion_result = new_page.insert_textbox(
                                                block_bbox,
                                                translated_text,
                                                fontname=fontname,
                                                fontsize=smaller_size,
                                                color=color,
                                                align=0,
                                                encoding=0  # Use default encoding
                                            )
                                            if insertion_result >= 0:
                                                break
                                
                            except Exception as e:
                                print(f"Error inserting text: {e}")
                                # Ultimate fallback: simple text insertion
                                try:
                                    new_page.insert_text(
                                        block_bbox.tl,
                                        translated_text,
                                        fontsize=12,
                                        color=(0, 0, 0),
                                        encoding=0  # Use default encoding
                                    )
                                except:
                                    # Last resort: insert original text
                                    new_page.insert_text(
                                        block_bbox.tl,
                                        original_text,
                                        fontsize=12,
                                        color=(0, 0, 0)
                                    )
            
            original_doc.close()
            
            # Save translated document to temporary file
            import tempfile
            import os
            
            temp_dir = Path(__file__).parent.parent / "uploads" / "temp"
            temp_dir.mkdir(parents=True, exist_ok=True)
            
            import uuid
            output_filename = f"translated_{uuid.uuid4().hex[:8]}.pdf"
            output_path = temp_dir / output_filename
            
            # Save with proper encoding options
            translated_doc.save(str(output_path), 
                              garbage=4,  # Maximum garbage collection
                              clean=True,  # Clean content streams
                              deflate=True)  # Compress
            translated_doc.close()
            
            return {
                'success': True,
                'output_path': str(output_path),
                'filename': output_filename,
                'translated_pages': len(page_numbers)
            }
            
        except Exception as e:
            return {'success': False, 'error': f'Basic PDF translation error: {str(e)}'}
    
    def _copy_page_images(self, source_page, target_page):
        """Copy all images from source page to target page."""
        try:
            # Get all images from the source page
            image_list = source_page.get_images()
            
            for img_index, img in enumerate(image_list):
                # Get image data
                xref = img[0]  # Image reference number
                try:
                    # Extract image
                    base_image = source_page.parent.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Get image rectangle
                    image_rects = source_page.get_image_rects(xref)
                    
                    if image_rects:
                        for rect in image_rects:
                            # Insert image into target page at the same location
                            target_page.insert_image(rect, stream=image_bytes)
                            
                except Exception as e:
                    print(f"Error copying image {img_index}: {e}")
                    continue
                    
        except Exception as e:
            print(f"Error in image copying: {e}")
    
    def _get_appropriate_font(self, target_language, font_info):
        """Get appropriate font based on target language and original font info."""
        original_font = font_info.get('font', 'Times-Roman').lower()
        is_bold = bool(font_info['flags'] & 2**4)
        is_italic = bool(font_info['flags'] & 2**6)
        
        # Language-specific font selection
        chinese_languages = ['traditional chinese', 'simplified chinese']
        asian_languages = ['japanese', 'korean', 'hindi']
        arabic_languages = ['arabic']
        
        if any(lang in target_language.lower() for lang in chinese_languages):
            # Use Chinese fonts for Chinese languages (both simplified and traditional)
            if is_bold and is_italic:
                return "china-ssi"  # or try "china-ss" 
            elif is_bold:
                return "china-ss"
            else:
                return "china-s"  # Chinese font that supports both simplified and traditional
        elif any(lang in target_language.lower() for lang in asian_languages):
            # Use CJK fonts for other Asian languages
            if is_bold and is_italic:
                return "china-ssi"  # CJK fonts often work for multiple Asian languages
            elif is_bold:
                return "china-ss"
            else:
                return "china-s"
        elif any(lang in target_language.lower() for lang in arabic_languages):
            # Use fonts that support Arabic
            return "noto-nastaliq"  # or "noto-arabic"
        else:
            # European languages - use Unicode-supporting fonts
            if 'times' in original_font:
                if is_bold and is_italic:
                    return "tibo"  # Times Bold Italic
                elif is_bold:
                    return "tibo"  # Times Bold
                elif is_italic:
                    return "tiro"  # Times Italic
                else:
                    return "tiro"  # Times Roman
            elif 'helvetica' in original_font or 'arial' in original_font:
                if is_bold and is_italic:
                    return "helv"  # Helvetica
                elif is_bold:
                    return "helv"
                elif is_italic:
                    return "helv"
                else:
                    return "helv"
            else:
                # Default to Helvetica for better Unicode support
                return "helv"
    
    def _convert_color(self, color_value):
        """Convert color from various formats to RGB tuple."""
        if isinstance(color_value, int):
            # Convert from integer to RGB
            r = (color_value >> 16) & 255
            g = (color_value >> 8) & 255
            b = color_value & 255
            return (r/255, g/255, b/255)
        elif isinstance(color_value, (list, tuple)) and len(color_value) >= 3:
            # Already RGB
            return tuple(color_value[:3])
        else:
            # Default black
            return (0, 0, 0)
    
    def _create_html_content(self, text, font_size, color, is_bold, is_italic):
        """Create HTML content for better Unicode text rendering."""
        # Convert color to hex
        if isinstance(color, tuple) and len(color) >= 3:
            r = int(color[0] * 255)
            g = int(color[1] * 255) 
            b = int(color[2] * 255)
            color_hex = f"#{r:02x}{g:02x}{b:02x}"
        else:
            color_hex = "#000000"
        
        # Check if text contains Chinese characters
        has_chinese = any('\u4e00' <= char <= '\u9fff' for char in text)
        
        # Choose appropriate font family based on content
        if has_chinese:
            font_family = "'SimSun', 'Microsoft YaHei', 'PingFang SC', 'Noto Sans CJK SC', 'Source Han Sans SC', sans-serif"
        else:
            font_family = "'DejaVu Sans', 'Noto Sans', Arial, sans-serif"
        
        # Build style
        style_parts = [
            f"font-size: {font_size}px",
            f"color: {color_hex}",
            f"font-family: {font_family}",
            "line-height: 1.2"
        ]
        
        if is_bold:
            style_parts.append("font-weight: bold")
        if is_italic:
            style_parts.append("font-style: italic")
            
        style = "; ".join(style_parts)
        
        # Escape HTML special characters but preserve Unicode
        import html
        escaped_text = html.escape(text, quote=False)
        
        return f'<div style="{style}">{escaped_text}</div>'
    
    def _extract_txt(self, file_path):
        """Extract text from TXT files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            return {'text': text}
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text = f.read()
                    return {'text': text, 'metadata': {'encoding': encoding}}
                except:
                    continue
            return {'error': 'Could not decode text file with any supported encoding'}
        except Exception as e:
            return {'error': f'Text extraction error: {str(e)}'}
    
    def _extract_docx(self, file_path):
        """Extract text from DOCX files."""
        try:
            from docx import Document
            doc = Document(file_path)
            
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(' | '.join(row_text))
            
            return {
                'text': '\n\n'.join(text_parts),
                'metadata': {
                    'paragraphs': len(doc.paragraphs),
                    'tables': len(doc.tables)
                }
            }
        except ImportError:
            return {'error': 'python-docx not installed. Run: pip install python-docx'}
        except Exception as e:
            return {'error': f'DOCX extraction error: {str(e)}'}
    
    def _extract_xlsx(self, file_path):
        """Extract text from XLSX files."""
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            text_parts = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None and str(cell).strip():
                            row_text.append(str(cell))
                    if row_text:
                        text_parts.append(' | '.join(row_text))
            
            return {
                'text': '\n'.join(text_parts),
                'metadata': {'sheets': len(workbook.sheetnames)}
            }
        except ImportError:
            return {'error': 'openpyxl not installed. Run: pip install openpyxl'}
        except Exception as e:
            return {'error': f'XLSX extraction error: {str(e)}'}
    
    def _extract_pptx(self, file_path):
        """Extract text from PPTX files."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            text_parts = []
            for i, slide in enumerate(prs.slides):
                text_parts.append(f"=== Slide {i+1} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            
            return {
                'text': '\n\n'.join(text_parts),
                'pages': len(prs.slides),
                'metadata': {'slides': len(prs.slides)}
            }
        except ImportError:
            return {'error': 'python-pptx not installed. Run: pip install python-pptx'}
        except Exception as e:
            return {'error': f'PPTX extraction error: {str(e)}'}
    
    def _extract_image(self, file_path):
        """Extract text from images using OCR."""
        try:
            import pytesseract
            from PIL import Image
            
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            
            return {
                'text': text,
                'metadata': {
                    'image_size': image.size,
                    'image_mode': image.mode
                }
            }
        except ImportError:
            return {'error': 'OCR libraries not installed. Run: pip install pytesseract pillow'}
        except Exception as e:
            return {'error': f'Image OCR error: {str(e)}'}
    
    def _extract_epub(self, file_path):
        """Extract text from EPUB files."""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            
            book = epub.read_epub(file_path)
            text_parts = []
            
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text = soup.get_text()
                    if text.strip():
                        text_parts.append(text.strip())
            
            return {
                'text': '\n\n--- Chapter Break ---\n\n'.join(text_parts),
                'metadata': {'chapters': len(text_parts)}
            }
        except ImportError:
            return {'error': 'EbookLib not installed. Run: pip install EbookLib beautifulsoup4'}
        except Exception as e:
            return {'error': f'EPUB extraction error: {str(e)}'}
    
    def _extract_svg(self, file_path):
        """Extract text from SVG files."""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            soup = BeautifulSoup(content, 'xml')
            
            # Extract text elements
            text_elements = soup.find_all('text')
            text_parts = [elem.get_text() for elem in text_elements if elem.get_text().strip()]
            
            return {
                'text': '\n'.join(text_parts),
                'metadata': {'text_elements': len(text_elements)}
            }
        except Exception as e:
            return {'error': f'SVG extraction error: {str(e)}'}
    
    def _extract_cbz(self, file_path):
        """Extract text from CBZ (Comic Book ZIP) files using OCR."""
        try:
            import pytesseract
            from PIL import Image
            import zipfile
            import io
            
            text_parts = []
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                image_files = [f for f in zip_file.namelist() 
                              if f.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp'))]
                image_files.sort()  # Sort to maintain order
                
                for image_file in image_files:
                    try:
                        image_data = zip_file.read(image_file)
                        image = Image.open(io.BytesIO(image_data))
                        text = pytesseract.image_to_string(image)
                        if text.strip():
                            text_parts.append(f"=== Page: {image_file} ===\n{text}")
                    except Exception as e:
                        text_parts.append(f"=== Error reading {image_file}: {str(e)} ===")
            
            return {
                'text': '\n\n'.join(text_parts),
                'pages': len(image_files),
                'metadata': {'comic_pages': len(image_files)}
            }
        except ImportError:
            return {'error': 'OCR libraries not installed. Run: pip install pytesseract pillow'}
        except Exception as e:
            return {'error': f'CBZ extraction error: {str(e)}'}
    
    def _extract_unsupported(self, file_path, extension):
        """Handle unsupported but recognized file types."""
        unsupported_info = {
            '.xps': 'XPS files require specialized libraries (not commonly available in Python)',
            '.mobi': 'MOBI files require Amazon\'s format libraries or conversion tools',
            '.fb2': 'FB2 files can be processed but require additional XML parsing libraries',
            '.hwp': 'HWP files require Korean office suite libraries (very limited Python support)'
        }
        
        return {
            'error': f'File type {extension} is recognized but not yet supported. {unsupported_info.get(extension, "")}'
        }

# Global instance
file_extractor = FileExtractor()
